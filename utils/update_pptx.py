import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def update_presentation():
    base_dir = "/Users/nickgklezakos/Documents/ppl-meta-code/docs/product/marketing/Digital signage"
    pptx_path = os.path.join(base_dir, "eyenet digital signage for technology retail shops.pptx")
    
    found_dir = None
    for d in os.listdir(base_dir):
        if d.startswith("signage photos"):
            found_dir = os.path.join(base_dir, d)
            break
            
    if not os.path.exists(pptx_path):
        print(f"Error: Presentation not found at {pptx_path}")
        return
    if not found_dir or not os.path.isdir(found_dir):
        print(f"Error: Image directory not found")
        return

    img_dir = found_dir
    prs = Presentation(pptx_path)
    
    img_exts = ('.png', '.jpg', '.jpeg', '.webp')
    img_files = sorted([f for f in os.listdir(img_dir) if f.lower().endswith(img_exts)])[:6]
    
    if not img_files:
        print(f"No images found in {img_dir}")
        return

    blank_slide_layout = prs.slide_layouts[6] 
    slide = prs.slides.add_slide(blank_slide_layout)
    
    color_indigo = RGBColor(0x3F, 0x51, 0xB5)
    color_border = RGBColor(0xDD, 0xDD, 0xDD)
    color_bg = RGBColor(0xFF, 0xFF, 0xFF)
    
    title_height = Inches(0.6)
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), title_height)
    title_tf = title_shape.text_frame
    title_tf.text = "Visual Gallery — AI Concept Frames"
    title_p = title_tf.paragraphs[0]
    title_p.font.size = Pt(28)
    title_p.font.bold = True
    title_p.font.color.rgb = color_indigo
    
    subtitle_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.4))
    subtitle_tf = subtitle_shape.text_frame
    subtitle_tf.text = "Representative persona visuals for Slides 2-9"
    subtitle_p = subtitle_tf.paragraphs[0]
    subtitle_p.font.size = Pt(14)
    subtitle_p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    rows, cols = 2, 3
    margin_top = Inches(1.4)
    margin_left = Inches(0.5)
    gutter = Inches(0.2)
    card_w = (Inches(10) - (2 * margin_left) - ((cols - 1) * gutter)) / cols
    card_h = (Inches(7.5) - margin_top - Inches(0.8) - ((rows - 1) * gutter)) / rows

    for i, img_name in enumerate(img_files):
        row = i // cols
        col = i % cols
        left = int(margin_left + col * (card_w + gutter))
        top = int(margin_top + row * (card_h + gutter))
        
        rect = slide.shapes.add_shape(1, left, top, int(card_w), int(card_h))
        rect.fill.solid()
        rect.fill.fore_color.rgb = color_bg
        line = rect.line
        line.color.rgb = color_border
        line.width = Pt(1)
        
        img_path = os.path.join(img_dir, img_name)
        inset = Inches(0.1)
        img_max_w = card_w - 2*inset
        img_max_h = card_h - Inches(0.5) - inset
        
        pic = slide.shapes.add_picture(img_path, left + int(inset), top + int(inset), width=int(img_max_w))
        if pic.height > img_max_h:
            ratio = img_max_h / pic.height
            pic.height = int(img_max_h)
            pic.width = int(pic.width * ratio)
            
        pic.left = int(left + (card_w - pic.width) / 2)
        pic.top = int(top + inset + (img_max_h - pic.height) / 2)

        cap_text = os.path.splitext(img_name)[0].replace('_', ' ').replace('-', ' ').title()
        cap_shape = slide.shapes.add_textbox(left, top + int(card_h - Inches(0.4)), int(card_w), int(Inches(0.3)))
        cap_tf = cap_shape.text_frame
        cap_tf.text = cap_text
        cap_p = cap_tf.paragraphs[0]
        cap_p.alignment = PP_ALIGN.CENTER
        cap_p.font.size = Pt(10)
        cap_p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    footer_y = int(Inches(7.1))
    footer_text = "Eyenet Vision  |  Technology Retail Pitch"
    footer_shape = slide.shapes.add_textbox(0, footer_y, int(Inches(10)), int(Inches(0.3)))
    footer_tf = footer_shape.text_frame
    footer_p = footer_tf.paragraphs[0]
    footer_p.text = footer_text
    footer_p.alignment = PP_ALIGN.CENTER
    footer_p.font.size = Pt(9)
    footer_p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

    prs.save(pptx_path)
    print(f"Images used: {len(img_files)}")
    print(f"Final slide count: {len(prs.slides)}")

if __name__ == "__main__":
    update_presentation()
