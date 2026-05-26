from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path("/Users/nickgklezakos/Documents/ppl-meta-code")
TEMPLATE_PPTX = ROOT / "docs/product/marketing/Digital signage/eyenet digital signage for technology retail shops.pptx"
OUTPUT_PPTX = ROOT / "docs/product/supportive/EYENET_VISION_FUNCTIONALITIES_PRESENTATION.pptx"
ASSETS_DIR = ROOT / "docs/product/business plan/assets"
SCREENSHOTS_DIR = ASSETS_DIR / "screenshots"

PRIMARY = RGBColor(0x1A, 0x27, 0x44)
ACCENT = RGBColor(0x25, 0x63, 0xEB)
ACCENT_LIGHT = RGBColor(0xDB, 0xEA, 0xFE)
TEXT = RGBColor(0x1E, 0x29, 0x3B)
TEXT_MUTED = RGBColor(0x64, 0x74, 0x8B)
SUCCESS = RGBColor(0x05, 0x96, 0x69)
CARD_BG = RGBColor(0xF8, 0xFA, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


CAPABILITY_SLIDES = [
    {
        "title": "How Eyenet Multiplies Human Operation Capacity",
        "headline": "Boost Moderate Cameras to Powerful Machine Vision Nodes",
        "quote": "Turn any camera into an intelligent sensor with real-time analytics and no specialized hardware burden.",
        "bullets": [
            "Turn IP, WiFi, and mobile cameras into intelligent machine-vision nodes.",
            "Run face detection, people counting, demographics, and route tracking on-device.",
            "Scale from a single lens to a facility-wide mesh without changing platform logic.",
        ],
        "image": ASSETS_DIR / "Platform Hero.png",
        "caption": "The Eyenet-vision Platform — Turning Every Camera into an Intelligent Sensor",
    },
    {
        "title": "Operational Coverage and Safety",
        "headline": "Track human activity, regulate access, and protect restricted zones",
        "quote": "Replace manual checks with always-on monitoring that counts, alerts, and documents events automatically.",
        "bullets": [
            "Real-time people tracking, attendance, route analysis, and dwell-time mapping.",
            "Automatic headcount, checkpoint oversight, and unauthorized-access alerts.",
            "Virtual perimeters, lone-worker monitoring, PPE checks, and incident documentation.",
        ],
        "image": ASSETS_DIR / "Eyenet security.jpg",
        "caption": "Advanced Security — Real-Time Monitoring and Intelligent Threat Detection",
    },
    {
        "title": "Audience-Aware Automation",
        "headline": "Drive digital signage and workflows from live vision events",
        "quote": "Cameras read the room and trigger content, notifications, and third-party actions in milliseconds.",
        "bullets": [
            "Switch playlists, promotions, and emergency messages based on age, gender, group size, and behavior.",
            "Configure per-camera triggers with configurable conditions and AND/OR logic.",
            "Measure impressions, dwell time, and engagement to prove signage ROI.",
        ],
        "image": ASSETS_DIR / "Digital signage.png",
        "caption": "Dynamic Digital Signage — Audience-Aware Content That Adapts in Real Time",
    },
    {
        "title": "Advanced Protection Use Cases",
        "headline": "Control entry points, restricted areas, VIP coverage, and remote operations",
        "quote": "One platform coordinates access control, perimeter alerts, VIP recognition, and remote supervision.",
        "bullets": [
            "Protect entrances, gates, and checkpoints with identity awareness and audit trails.",
            "Stream person-of-interest alerts to security eyewear and mobile devices.",
            "Manage every camera, trigger, and dashboard remotely from web or mobile with no cloud video dependency.",
        ],
        "image": ASSETS_DIR / "VIP Protection.png",
        "caption": "VIP Protection — Discreet Recognition and Priority Alert Systems",
    },
    {
        "title": "Instant Deployment",
        "headline": "Get started quickly on everyday hardware",
        "quote": "Deploy with standard compute, no forced GPU servers, and no per-camera licensing ceiling.",
        "bullets": [
            "Start with 16 GB RAM and a modern CPU using cameras you already own.",
            "Bring the first detection online in minutes from IP, WiFi, or phone-based cameras.",
            "Expand from one room to multi-site estates without redesigning the architecture.",
        ],
        "image": ASSETS_DIR / "Instant Deployment.png",
        "caption": "Instant Deployment — Up and Running in Minutes, Not Months",
    },
]

DIFFERENTIATORS = [
    "Privacy-first on-premises processing and storage.",
    "Offline-capable operation without internet dependency.",
    "Scales from one camera to enterprise-wide deployments.",
    "Secure by design with encryption, access control, and VPN-ready connectivity.",
    "Integration-ready with REST APIs, webhooks, and third-party connectors.",
    "Event-driven intelligence with per-camera automation and multi-action orchestration.",
]

MODULE_GROUPS = [
    {
        "title": "Core Platform Functionalities",
        "quote": "The platform combines computer vision, hardware integration, triggers, media handling, and communications into one operating layer.",
        "columns": [
            (
                "Computer Vision, Cameras, Triggers",
                [
                    "People analytics: tracking, counting, attendance, route analysis, and demographics.",
                    "Multi-source camera support across IP, WiFi, and mobile soft cameras.",
                    "Per-camera trigger configuration with real-time event detection and complex condition logic.",
                    "Automated actions across push, in-app, email, SMS, signage, and webhooks.",
                ],
            ),
            (
                "Video, Data, and Communications",
                [
                    "Local video storage and processing with no mandatory cloud services.",
                    "Intelligent communications with audit logging and delivery tracking.",
                    "Real-time dashboards, historical analytics, behavioral patterns, and reporting.",
                    "Operational resilience with internet optional except for selected remote notifications.",
                ],
            ),
        ],
    },
    {
        "title": "Architecture, Security, and Scale",
        "quote": "Eyenet is engineered as an edge-capable microservices platform that preserves control while scaling across sites and stakeholders.",
        "columns": [
            (
                "Deployment and Governance",
                [
                    "On-premises and edge deployment with modular microservices and automatic service discovery.",
                    "Role-based access, file-level encryption, private VPN support, and full data sovereignty.",
                    "Remote diagnostics, centralized logging, audit trails, and trigger execution debugging.",
                ],
            ),
            (
                "Scalability and Extensibility",
                [
                    "Flexible scaling from homes to industrial sites with hardware-based performance growth.",
                    "Unlimited camera support subject to available hardware capacity.",
                    "Extensible integrations through REST APIs, webhooks, custom connectors, and signage modules.",
                ],
            ),
        ],
    },
    {
        "title": "Technical Architecture Highlights",
        "quote": "Trigger evaluation, orchestration, and context awareness are first-class capabilities rather than bolt-on automation.",
        "columns": [
            (
                "Trigger and Action System",
                [
                    "Granular per-camera rule configuration.",
                    "Instant event evaluation against configured conditions.",
                    "Reliable action delivery with retry logic and failure tracking.",
                    "Complete audit trail for activations and executions.",
                ],
            ),
            (
                "Context and Signage Intelligence",
                [
                    "Time, day, occupancy, demographics, velocity, direction, and trajectory-aware triggers.",
                    "Historical comparison and multi-camera correlation across zones.",
                    "Dynamic signage content switching, demographic targeting, and playlist automation.",
                ],
            ),
        ],
    },
]

SCREENSHOT_SLIDES = [
    {
        "title": "Eyenet in Action — Operations and Analytics",
        "quote": "Representative screenshots from the live platform UX included in the functionalities document.",
        "items": [
            ("main-dashboard.png", "Main Dashboard"),
            ("anaytics-filters.png", "Analytics Filters"),
            ("settings.png", "System Settings"),
            ("analytics-01.png", "Analytics Dashboard"),
            ("camera-analytics-03.png", "Detection Analytics"),
            ("camera-list.png", "Camera Management"),
        ],
    },
    {
        "title": "Eyenet in Action — Automation and Media",
        "quote": "Workflow and endpoint views covering trigger orchestration, signage, media, and playback operations.",
        "items": [
            ("triggers-list.png", "Triggers Management"),
            ("actions-edit.png", "Actions Editor"),
            ("digital-signage-01.png", "Digital Signage"),
            ("mvr-groups-details.png", "MVR Group Details"),
            ("upload-media.png", "Media Upload"),
            ("signage-android-app.jpg", "Signage Android App"),
        ],
    },
]


def remove_all_slides(prs: Presentation) -> None:
    slide_ids = list(prs.slides._sldIdLst)
    for slide_id in slide_ids:
        rel_id = slide_id.rId
        prs.part.drop_rel(rel_id)
        prs.slides._sldIdLst.remove(slide_id)


def add_footer(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.4), Inches(7.05), Inches(9.2), Inches(0.25))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(9)
    p.font.color.rgb = TEXT_MUTED


def add_quote(slide, text: str, top: float) -> None:
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(top), Inches(5.0), Inches(0.72))
    box.fill.solid()
    box.fill.fore_color.rgb = ACCENT_LIGHT
    box.line.color.rgb = ACCENT_LIGHT
    tf = box.text_frame
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.18)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = f'"{text}"'
    p.font.size = Pt(15)
    p.font.italic = True
    p.font.color.rgb = PRIMARY


def style_title(shape, size: int = 24) -> None:
    p = shape.text_frame.paragraphs[0]
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = PRIMARY


def add_bullets_box(slide, bullets, left, top, width, height, font_size=16) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = Inches(0.05)
    for index, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.bullet = True
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT
        p.space_after = Pt(9)
        p.line_spacing = 1.15


def add_image_card(slide, image_path: Path, caption: str, left: float, top: float, width: float, height: float) -> None:
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = PRIMARY
    card.line.color.rgb = PRIMARY

    inset = Inches(0.08)
    caption_h = Inches(0.52)
    image_area_h = height - caption_h - inset * 2
    slide.shapes.add_picture(str(image_path), left + inset, top + inset, width=width - inset * 2, height=image_area_h)

    caption_box = slide.shapes.add_textbox(left + Inches(0.15), top + height - caption_h + Inches(0.05), width - Inches(0.3), caption_h - Inches(0.08))
    tf = caption_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = caption
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE


def build_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Eyenet-vision Platform"
    style_title(slide.shapes.title, size=28)

    subtitle = slide.placeholders[1].text_frame
    subtitle.clear()
    p1 = subtitle.paragraphs[0]
    p1.text = "Core Functionalities & Features"
    p1.font.size = Pt(22)
    p1.font.bold = True
    p1.font.color.rgb = ACCENT
    p2 = subtitle.add_paragraph()
    p2.text = "Version 2.23.1  |  January 15, 2026  |  Confidential"
    p2.font.size = Pt(14)
    p2.font.color.rgb = TEXT_MUTED

    slide.shapes.add_picture(str(ASSETS_DIR / "eyenet-logo.png"), Inches(0.55), Inches(0.35), height=Inches(0.75))
    add_quote(slide, "Turning ordinary cameras into an intelligent, always-on assistant for safer, smarter, more efficient operations.", 1.75)
    add_bullets_box(
        slide,
        [
            "Counts visitors, detects unusual events, and drives automated responses in real time.",
            "Runs locally on owner-controlled hardware with privacy-first, offline-capable operation.",
            "Scales from a single camera to multi-site estates without per-camera licensing ceilings.",
        ],
        Inches(0.7),
        Inches(2.8),
        Inches(4.7),
        Inches(2.5),
        font_size=17,
    )
    add_image_card(
        slide,
        ASSETS_DIR / "Platform Hero.png",
        "The Eyenet-vision Platform — Turning Every Camera into an Intelligent Sensor",
        Inches(5.7),
        Inches(1.65),
        Inches(3.8),
        Inches(4.8),
    )
    add_footer(slide, "Eyenet Vision  |  Functionalities Presentation")


def build_capability_slide(prs: Presentation, spec: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(8.7), Inches(0.45))
    title.text_frame.text = spec["title"]
    style_title(title, size=24)

    headline = slide.shapes.add_textbox(Inches(0.55), Inches(0.95), Inches(5.0), Inches(0.55))
    htf = headline.text_frame
    p = htf.paragraphs[0]
    p.text = spec["headline"]
    p.font.size = Pt(21)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    add_quote(slide, spec["quote"], 1.55)
    add_bullets_box(slide, spec["bullets"], Inches(0.7), Inches(2.55), Inches(4.8), Inches(2.8), font_size=16)
    add_image_card(slide, spec["image"], spec["caption"], Inches(5.65), Inches(1.35), Inches(3.75), Inches(4.95))
    add_footer(slide, "Eyenet Vision  |  Functionalities Presentation")


def build_differentiators_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(8.5), Inches(0.45))
    title.text_frame.text = "Key Differentiators"
    style_title(title, size=24)
    add_quote(slide, "A privacy-first machine-vision platform designed for owner control, flexible automation, and practical scale.", 0.95)

    card_width = Inches(4.15)
    card_height = Inches(0.72)
    start_left = Inches(0.6)
    start_top = Inches(2.0)
    gutter_x = Inches(0.3)
    gutter_y = Inches(0.22)

    for idx, item in enumerate(DIFFERENTIATORS):
        row = idx // 2
        col = idx % 2
        left = start_left + col * (card_width + gutter_x)
        top = start_top + row * (card_height + gutter_y)
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, card_width, card_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.color.rgb = CARD_BG

        icon = slide.shapes.add_textbox(left + Inches(0.12), top + Inches(0.12), Inches(0.3), Inches(0.3))
        ip = icon.text_frame.paragraphs[0]
        ip.text = "✓"
        ip.font.size = Pt(18)
        ip.font.bold = True
        ip.font.color.rgb = SUCCESS

        text_box = slide.shapes.add_textbox(left + Inches(0.45), top + Inches(0.09), card_width - Inches(0.58), Inches(0.45))
        tp = text_box.text_frame.paragraphs[0]
        tp.text = item
        tp.font.size = Pt(15)
        tp.font.color.rgb = TEXT

    add_image_card(slide, ASSETS_DIR / "Access control.png", "Intelligent Access Control — Seamless Entry Management Powered by Vision AI", Inches(6.15), Inches(4.85), Inches(3.15), Inches(1.9))
    add_footer(slide, "Eyenet Vision  |  Functionalities Presentation")


def build_columns_slide(prs: Presentation, spec: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(8.6), Inches(0.45))
    title.text_frame.text = spec["title"]
    style_title(title, size=24)
    add_quote(slide, spec["quote"], 0.95)

    positions = [Inches(0.6), Inches(5.1)]
    for idx, (heading, bullets) in enumerate(spec["columns"]):
        left = positions[idx]
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, Inches(1.95), Inches(4.0), Inches(4.45))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BG

        hb = slide.shapes.add_textbox(left + Inches(0.2), Inches(2.1), Inches(3.6), Inches(0.4))
        hp = hb.text_frame.paragraphs[0]
        hp.text = heading
        hp.font.size = Pt(18)
        hp.font.bold = True
        hp.font.color.rgb = ACCENT

        add_bullets_box(slide, bullets, left + Inches(0.2), Inches(2.55), Inches(3.6), Inches(3.5), font_size=15)

    add_footer(slide, "Eyenet Vision  |  Functionalities Presentation")


def build_gallery_slide(prs: Presentation, spec: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(8.8), Inches(0.45))
    title.text_frame.text = spec["title"]
    style_title(title, size=24)
    add_quote(slide, spec["quote"], 0.95)

    rows = 2
    cols = 3
    card_w = Inches(2.78)
    card_h = Inches(2.1)
    start_left = Inches(0.6)
    start_top = Inches(2.0)
    gutter_x = Inches(0.2)
    gutter_y = Inches(0.25)

    for idx, (filename, label) in enumerate(spec["items"]):
        row = idx // cols
        col = idx % cols
        left = start_left + col * (card_w + gutter_x)
        top = start_top + row * (card_h + gutter_y)

        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = CARD_BG

        slide.shapes.add_picture(str(SCREENSHOTS_DIR / filename), left + Inches(0.07), top + Inches(0.07), width=card_w - Inches(0.14), height=Inches(1.45))
        label_box = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(1.58), card_w - Inches(0.2), Inches(0.3))
        lp = label_box.text_frame.paragraphs[0]
        lp.text = label
        lp.font.size = Pt(11)
        lp.font.bold = True
        lp.font.color.rgb = PRIMARY
        lp.alignment = PP_ALIGN.CENTER

    add_footer(slide, "Eyenet Vision  |  Functionalities Presentation")


def generate() -> None:
    prs = Presentation(str(TEMPLATE_PPTX))
    remove_all_slides(prs)

    build_title_slide(prs)
    for spec in CAPABILITY_SLIDES:
        build_capability_slide(prs, spec)
    build_differentiators_slide(prs)
    for spec in MODULE_GROUPS:
        build_columns_slide(prs, spec)
    for spec in SCREENSHOT_SLIDES:
        build_gallery_slide(prs, spec)

    prs.save(str(OUTPUT_PPTX))
    print(f"Created {OUTPUT_PPTX}")
    print(f"Slide count: {len(prs.slides)}")


if __name__ == "__main__":
    generate()